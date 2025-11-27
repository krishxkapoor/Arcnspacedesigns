from pptx import Presentation
from pptx.util import Inches, Pt
import os

TEMPLATES_DIR = "templates"
os.makedirs(TEMPLATES_DIR, exist_ok=True)

def create_template(filename, title):
    prs = Presentation()
    slide_layout = prs.slide_layouts[0] # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    title_placeholder = slide.shapes.title
    subtitle_placeholder = slide.placeholders[1]
    
    title_placeholder.text = title
    subtitle_placeholder.text = "Valuation Report Template"
    
    # Add a second slide for details
    bullet_layout = prs.slide_layouts[1]
    slide2 = prs.slides.add_slide(bullet_layout)
    shapes = slide2.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "Property Details"
    tf = body_shape.text_frame
    tf.text = "Client Name: {{client_name}}"
    
    p = tf.add_paragraph()
    p.text = "Date: {{date}}"
    
    p = tf.add_paragraph()
    p.text = "Purpose: {{purpose}}"
    
    prs.save(os.path.join(TEMPLATES_DIR, filename))
    print(f"Created {filename}")

if __name__ == "__main__":
    create_template("Agriculture.pptx", "Agricultural Valuation")
    create_template("Commercial.pptx", "Commercial Valuation")
    create_template("Residential.pptx", "Residential Valuation")
