"""
Script to inspect PowerPoint template structure
"""
from pptx import Presentation
import sys

def inspect_template(template_path):
    """Inspect a PowerPoint template to find text placeholders"""
    prs = Presentation(template_path)
    
    print(f"\n{'='*60}")
    print(f"Template: {template_path}")
    print(f"{'='*60}\n")
    
    for slide_idx, slide in enumerate(prs.slides):
        print(f"Slide {slide_idx + 1}:")
        print("-" * 40)
        
        for shape_idx, shape in enumerate(slide.shapes):
            if hasattr(shape, "text"):
                print(f"  Shape {shape_idx}: {shape.text[:50] if shape.text else '(empty)'}")
                if hasattr(shape, "text_frame"):
                    for para_idx, paragraph in enumerate(shape.text_frame.paragraphs):
                        for run_idx, run in enumerate(paragraph.runs):
                            if run.text.strip():
                                print(f"    - Text: '{run.text}'")
                                font_name = run.font.name if run.font.name else "Default"
                                font_size = run.font.size if run.font.size else "Default"
                                print(f"      Font: {font_name}, Size: {font_size}")
                                try:
                                    if run.font.color and hasattr(run.font.color, 'rgb') and run.font.color.rgb:
                                        print(f"      Color: RGB{run.font.color.rgb}")
                                except:
                                    pass
        print()

if __name__ == "__main__":
    templates = [
        "templates/Agriculture.pptx",
        "templates/Commercial.pptx",
        "templates/Residential.pptx"
    ]
    
    for template in templates:
        try:
            inspect_template(template)
        except Exception as e:
            print(f"Error inspecting {template}: {e}")
